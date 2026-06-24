#!/usr/bin/env python3
"""
build_slides.py
Commodities-desk slide deck (python-pptx) for the refinery-outage analysis.

Priority-2 deliverable. Mirrors the workbook's charts as crisp embedded images
(rendered by charts.py) in a clean 16:9 layout, with data-driven insight
captions that refresh with the data. Self-contained: charts are embedded, not
linked.

Usage:
    python build_slides.py                       # uses INPUT_PATH
    python build_slides.py path/to/export.xlsx --out outage_deck.pptx
"""
import argparse
import os
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import engine
import charts

INPUT_PATH = "rEFINERY oUTAGES.xlsx"
OUT_PATH = "outage_deck.pptx"

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x54, 0x96)
RED = RGBColor(0xC0, 0x00, 0x00)
GOLD = RGBColor(0xBF, 0x90, 0x00)
GREEN = RGBColor(0x54, 0x82, 0x35)
LT_BLUE = RGBColor(0xD6, 0xE0, 0xF0)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

FONT = "Arial"
SW, SH = Inches(13.333), Inches(7.5)


def kbd(x):
    return f"{x:,.0f}"


class Deck:
    def __init__(self, ctx, assets):
        self.ctx = ctx
        self.assets = assets
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.blank = self.prs.slide_layouts[6]

    # ------------------------------------------------------------- primitives
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, s, l, t, w, h, color, line=None):
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(1)
        shp.shadow.inherit = False
        return shp

    def _text(self, s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              space_after=4):
        """runs: list of (text, size, bold, color, italic)."""
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Pt(2)
        tf.margin_top = tf.margin_bottom = Pt(2)
        first = True
        for item in runs:
            text, size, bold, color = item[0], item[1], item[2], item[3]
            italic = item[4] if len(item) > 4 else False
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space_after)
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = FONT
            r.font.color.rgb = color
        return tb

    def _titlebar(self, s, title, sub=None):
        self._rect(s, 0, 0, SW, Inches(0.9), NAVY)
        self._text(s, Inches(0.45), Inches(0.12), Inches(12.4), Inches(0.5),
                   [(title, 24, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
        if sub:
            self._text(s, Inches(0.47), Inches(0.6), Inches(12.4), Inches(0.28),
                       [(sub, 11, False, LT_BLUE)])
        # gold accent under bar
        self._rect(s, 0, Inches(0.9), SW, Pt(3), GOLD)

    def _footer(self, s, page):
        d = self.ctx["diag"]
        y0, y1 = d["years"]
        self._text(s, Inches(0.45), Inches(7.12), Inches(9), Inches(0.3),
                   [(f"Refinery Outage Analytics  |  Source: Snowflake export "
                     f"({d['rows']:,} rows, {y0}-{y1})  |  Primary metric: capacity offline (kbd)",
                     8, False, GRAY)])
        self._text(s, Inches(12.4), Inches(7.12), Inches(0.6), Inches(0.3),
                   [(str(page), 8, False, GRAY)], align=PP_ALIGN.RIGHT)

    def _pic_fit(self, s, path, l, t, max_w, max_h, center=True):
        """Insert a picture scaled to fit (max_w, max_h), preserving aspect."""
        pic = s.shapes.add_picture(path, l, t, width=max_w)
        if pic.height > max_h:
            scale = max_h / pic.height
            pic.height = int(pic.height * scale)
            pic.width = int(pic.width * scale)
        if center:
            pic.left = int(l + (max_w - pic.width) / 2)
        return pic

    def _chart_slide(self, title, sub, img, insight, page, img_h=Inches(5.0)):
        s = self._slide()
        self._titlebar(s, title, sub)
        self._pic_fit(s, img, Inches(0.5), Inches(1.15), Inches(12.3), img_h)
        if insight:
            self._rect(s, Inches(0.5), Inches(6.35), Inches(12.33), Inches(0.62), LT_GRAY)
            self._rect(s, Inches(0.5), Inches(6.35), Inches(0.09), Inches(0.62), GOLD)
            self._text(s, Inches(0.75), Inches(6.4), Inches(11.9), Inches(0.55),
                       [("Takeaway  ", 10, True, NAVY), (insight, 10, False, RGBColor(0x40, 0x40, 0x40))],
                       anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, page)
        return s

    # ------------------------------------------------------------- slides
    def title_slide(self):
        s = self._slide()
        self._rect(s, 0, 0, SW, SH, NAVY)
        self._rect(s, 0, Inches(3.05), SW, Pt(3), GOLD)
        d = self.ctx["diag"]
        y0, y1 = d["years"]
        self._text(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.0),
                   [("Refinery Outage Analytics", 44, True, WHITE)])
        self._text(s, Inches(0.82), Inches(3.2), Inches(11.7), Inches(0.6),
                   [("Capacity Offline (kbd) - Planned & Unplanned - 2027 Scenario & Sensitivity",
                     18, False, LT_BLUE)])
        s_ = self.ctx["summary"]
        ly = max(y for y in s_.index if y not in engine.PARTIAL_YEARS and s_.loc[y, "Unplanned"] > 0)
        self._text(s, Inches(0.82), Inches(4.5), Inches(11.7), Inches(2.0),
                   [(f"Data vintage:  {d['rows']:,} rows  |  {y0}-{y1}  |  "
                     f"{d['events_distinct']:,} distinct outages", 13, False, WHITE),
                    (f"Latest full year ({ly}):  {kbd(s_.loc[ly,'Total'])} kbd total offline, "
                     f"{kbd(s_.loc[ly,'Unplanned'])} kbd unplanned ({s_.loc[ly,'Unpl%']:.0%})",
                     13, False, WHITE),
                    ("2026 & 2027 are partial / planned-only - shown for context, not as finals.",
                     11, False, RGBColor(0xF0, 0xC0, 0xC0), True)], space_after=10)

    def exec_summary(self):
        s = self._slide()
        self._titlebar(s, "Executive Summary", "At-a-glance metrics and what matters")
        sm = self.ctx["summary"]
        ly = max(y for y in sm.index if y not in engine.PARTIAL_YEARS and sm.loc[y, "Unplanned"] > 0)
        padd_un = self.ctx["padd_unplanned"]
        top_padd = padd_un[ly].idxmax()
        share = self.ctx["padd_share"]
        sc = self.ctx["scenario"]
        tiles = [
            (f"Total Offline (FY{ly})", kbd(sm.loc[ly, "Total"]), "kbd"),
            (f"Unplanned (FY{ly})", kbd(sm.loc[ly, "Unplanned"]), "kbd"),
            ("Unplanned %", f"{sm.loc[ly,'Unpl%']:.0%}", "of total"),
            ("Distinct Outages", kbd(sm.loc[ly, "Events"]), f"FY{ly}"),
            ("Top PADD", top_padd, "by unplanned"),
        ]
        n = len(tiles)
        gap = Inches(0.2)
        tw = (SW - Inches(0.9) - gap * (n - 1)) / n
        x = Inches(0.45)
        for lab, val, sub in tiles:
            self._rect(s, x, Inches(1.25), tw, Inches(0.42), BLUE)
            self._text(s, x, Inches(1.27), tw, Inches(0.4),
                       [(lab, 10.5, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._rect(s, x, Inches(1.67), tw, Inches(1.1), LT_BLUE)
            self._text(s, x, Inches(1.72), tw, Inches(0.7),
                       [(val, 26, True, NAVY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._text(s, x, Inches(2.4), tw, Inches(0.3),
                       [(sub, 9, False, GRAY, True)], align=PP_ALIGN.CENTER)
            x += tw + gap

        # takeaways
        yoy = sm.loc[ly, "YoY%"]
        peak_season = "February (winter freeze) and September-October (turnaround)"
        bullets = [
            (f"{top_padd} dominates unplanned risk at ~{share[top_padd]:.0%} of US unplanned "
             f"capacity offline over the baseline window; PADD 2 and PADD 5 are the next largest.",),
            (f"FY{ly} total offline of {kbd(sm.loc[ly,'Total'])} kbd was {yoy:+.0%} YoY, with the "
             f"planned/unplanned split near 50/50 ({sm.loc[ly,'Unpl%']:.0%} unplanned).",),
            (f"Unplanned capacity offline is highly seasonal - peaks in {peak_season}, "
             f"with a summer trough.",),
            (f"2027 is planned-only in the data ({kbd(sc['planned_2027'])} kbd booked). The scenario "
             f"model adds a baseline-driven unplanned forecast of ~{kbd(sc['annual_unplanned'])} kbd "
             f"-> implied total ~{kbd(sc['implied_total'])} kbd.",),
        ]
        self._text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.35),
                   [("Key Takeaways", 14, True, NAVY)])
        ty = Inches(3.45)
        for b in bullets:
            self._rect(s, Inches(0.55), ty + Inches(0.06), Inches(0.12), Inches(0.12), GOLD)
            self._text(s, Inches(0.85), ty, Inches(11.9), Inches(0.8),
                       [(b[0], 12.5, False, RGBColor(0x33, 0x33, 0x33))])
            ty += Inches(0.82)
        self._footer(s, 2)

    def build(self):
        a = self.assets
        self.title_slide()
        self.exec_summary()
        sm = self.ctx["summary"]
        ly = max(y for y in sm.index if y not in engine.PARTIAL_YEARS and sm.loc[y, "Unplanned"] > 0)
        peak_yr = max((y for y in sm.index if 2014 <= y <= 2025), key=lambda y: sm.loc[y, "Total"])
        share = self.ctx["padd_share"]
        sc = self.ctx["scenario"]

        self._chart_slide(
            "Annual Trend", "Capacity offline by year, planned + unplanned",
            a["annual"],
            f"Offline capacity peaked in {peak_yr} (COVID / Winter Storm Uri years 2020-21 are "
            f"outliers and excluded from forecast baselines). Recent years run a roughly even "
            f"planned/unplanned split.", 3)

        # PADD: two images side by side
        s = self._slide()
        self._titlebar(s, "Unplanned by PADD", "Where the unplanned risk concentrates")
        self._pic_fit(s, a["padd_clustered"], Inches(0.45), Inches(1.2), Inches(7.0), Inches(4.9), center=False)
        self._pic_fit(s, a["padd_donut"], Inches(7.7), Inches(1.2), Inches(5.3), Inches(4.9), center=False)
        self._rect(s, Inches(0.5), Inches(6.35), Inches(12.33), Inches(0.62), LT_GRAY)
        self._rect(s, Inches(0.5), Inches(6.35), Inches(0.09), Inches(0.62), GOLD)
        self._text(s, Inches(0.75), Inches(6.4), Inches(11.9), Inches(0.55),
                   [("Takeaway  ", 10, True, NAVY),
                    (f"PADD 3 (Gulf Coast) carries ~{share['PADD 3']:.0%} of US unplanned capacity "
                     f"offline; PADD 5 spiked recently on California events.", 10, False,
                     RGBColor(0x40, 0x40, 0x40))], anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, 4)

        self._chart_slide(
            "Seasonality", "Unplanned offline by calendar month",
            a["seasonality"],
            "Unplanned outages cluster in late winter (freeze events) and the autumn turnaround "
            "window, with a summer trough - the shape that drives the 2027 scenario.", 5)

        self._chart_slide(
            "PADD 3 - Reference View", "2026 plan + unplanned vs 2023-2025 totals and 2027 plan",
            a["padd3"],
            "Gold/orange columns are 2026 booked plan and unplanned; lines are prior-year total "
            "offline; the dashed green line is 2027 booked plan. Magnitudes share one axis (honest "
            "scale).", 6)

        self._chart_slide(
            "PADD 1, 2, 4 & 5", "Same view across the remaining PADDs",
            a["padd_sm"],
            "PADD 1 is genuinely small; PADD 2 and PADD 5 show the largest swings outside the Gulf "
            "Coast. 2027 plan (dashed green) is the only forward-booked series.", 7)

        self._chart_slide(
            "Unit Categories", "Capacity offline by unit category (all years)",
            a["units"],
            "Crude (atmospheric/vacuum distillation), FCC and reforming units account for the bulk "
            "of offline capacity - and the bulk of gasoline-relevant exposure.", 8)

        self._chart_slide(
            "2027 Unplanned Scenario", "Driver-based forecast vs planned and recent actuals",
            a["scenario"],
            f"Default scenario (window {sc['window']}, 0% growth, 1.0x multiplier): "
            f"~{kbd(sc['annual_unplanned'])} kbd unplanned on top of {kbd(sc['planned_2027'])} kbd "
            f"booked plan -> ~{kbd(sc['implied_total'])} kbd implied total. Fully tunable in the "
            f"workbook.", 9)

        # Sensitivity: heatmap + tornado
        s = self._slide()
        self._titlebar(s, "Sensitivity & Risk", "How the 2027 unplanned forecast flexes")
        self._pic_fit(s, a["heatmap"], Inches(0.45), Inches(1.2), Inches(6.7), Inches(4.9), center=False)
        self._pic_fit(s, a["tornado"], Inches(7.35), Inches(1.45), Inches(5.7), Inches(4.4), center=False)
        self._rect(s, Inches(0.5), Inches(6.35), Inches(12.33), Inches(0.62), LT_GRAY)
        self._rect(s, Inches(0.5), Inches(6.35), Inches(0.09), Inches(0.62), GOLD)
        torn_top = self.ctx["tornado"][0]["driver"].split(" (")[0]
        self._text(s, Inches(0.75), Inches(6.4), Inches(11.9), Inches(0.55),
                   [("Takeaway  ", 10, True, NAVY),
                    (f"The unplanned-rate multiplier is the dominant driver ({torn_top}); the "
                     f"heatmap base case ({kbd(sc['annual_unplanned'])} kbd) is outlined.", 10,
                     False, RGBColor(0x40, 0x40, 0x40))], anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, 10)

        # closing / method
        s = self._slide()
        self._titlebar(s, "Method & Caveats", "Read before use")
        notes = [
            "Primary metric is CAP_OFFLINE_ADJUSTED_KBD (offline capacity, kbd, all units). Mogas is "
            "a secondary overlay only.",
            "Outage type is binary {PLANNED, UNPLANNED}; UNKNOWN folds into UNPLANNED per desk rule.",
            "2027 is planned-only - any Plan+Unplanned or Unplanned comparison vs 2027 is shown n/a; "
            "only Planned is comparable. Unplanned-2027 is a modeled scenario.",
            "2020-2021 (COVID / Winter Storm Uri) are excluded from forecast baselines by default.",
            "Scenario = baseline(window) x (1+growth) x multiplier + one-off(stress month); the Excel "
            "workbook recomputes it live from dropdown inputs.",
            "All numbers refresh from the source export - re-run the build to update workbook, deck "
            "and dashboard together.",
        ]
        ty = Inches(1.4)
        for nlabel in notes:
            self._rect(s, Inches(0.55), ty + Inches(0.07), Inches(0.12), Inches(0.12), GOLD)
            self._text(s, Inches(0.85), ty, Inches(11.9), Inches(0.8),
                       [(nlabel, 12, False, RGBColor(0x33, 0x33, 0x33))])
            ty += Inches(0.78)
        self._footer(s, 11)

    def save(self, path):
        self.prs.save(path)


def main():
    ap = argparse.ArgumentParser(description="Refinery-outage slide deck")
    ap.add_argument("excel", nargs="?", default=INPUT_PATH, help="path to the outage .xlsx export")
    ap.add_argument("--out", default=OUT_PATH, help="output .pptx path")
    args = ap.parse_args()

    print(f"Loading {args.excel.strip()} ...")
    ctx = engine.build_context(args.excel)
    with tempfile.TemporaryDirectory() as tmp:
        print("Rendering charts ...")
        assets = charts.render_all(ctx, tmp)
        print(f"Building deck -> {args.out}")
        deck = Deck(ctx, assets)
        deck.build()
        deck.save(args.out)
    print("Done.")


if __name__ == "__main__":
    sys.exit(main())
