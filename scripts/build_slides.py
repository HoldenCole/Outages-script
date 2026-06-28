#!/usr/bin/env python3
"""
build_slides.py
Trading-desk deck (python-pptx), evenly split gasoline / distillate. Slides are
chart-forward; the talk-track is in the per-slide footnotes.
    1. Total 2027 outages by unit  (all four focus units, confirmed vs not-booked)
    2. What's driving the numbers  (the biggest individual outages, by PADD)
    3. Month-over-month change     (what moved this month, by region & focus unit)
    4. H1 planned per unit & month (2025 vs 2026 vs 2027, like-for-like)
    5. Gasoline complex by PADD    (CDU + FCC)
    6. Distillate complex by PADD  (CDU + hydrocracker)
    7. Unplanned recent context    (grounds the 2027 scenario)
    8. 2027 unplanned scenario     (monthly risk on top of the booked plan)
    9. What it means for the market (gasoline & distillate, thin inventories)

Everything is per-unit capacity offline (never a summed "total"), 2027-forward.
2027 completeness is asymmetric: only ExxonMobil gave a full-year plan (verified
against their corporate schedule); every other operator is H1-confirmed only, so
the non-Exxon H2 is flagged as not-yet-booked throughout.

Brand mark: set BRAND_LOGO to a logo image path; otherwise BRAND_TEXT is used.

Usage:
    python build_slides.py                       # uses INPUT_PATH
    python build_slides.py path/to/export.xlsx --out outage_deck.pptx
"""
import argparse
import math
import os
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import engine
import charts

_ROOT = Path(__file__).resolve().parent.parent
INPUT_PATH = str(_ROOT / "data" / "Golden_Record_Snowflake.xlsx")
OUT_PATH = str(_ROOT / "output" / "outage_deck.pptx")

BRAND_TEXT = "Products Trading"
BRAND_LOGO = None

FY = engine.FOCUS_YEAR              # forward outlook year (current year + 1); rolls with the data
Y0 = engine.START_YEAR             # 2023

NAVY = RGBColor(0x1F, 0x38, 0x64)
NAVY2 = RGBColor(0x2E, 0x54, 0x96)
RED = RGBColor(0xC0, 0x00, 0x00)
GOLD = RGBColor(0xBF, 0x90, 0x00)
GREEN = RGBColor(0x54, 0x82, 0x35)
LT_BLUE = RGBColor(0xD6, 0xE0, 0xF0)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
INK = RGBColor(0x26, 0x2B, 0x33)

FONT = "Arial"
SW, SH = Inches(13.333), Inches(7.5)


def kbd(x):
    return f"{x:,.0f}"


class Deck:
    def __init__(self, ctx, assets, asof="June 26th, 2026"):
        self.ctx = ctx
        self.a = assets
        self.asof = asof
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.blank = self.prs.slide_layouts[6]
        self.page = 0

    # ----------------------------------------------------------------- primitives
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, s, l, t, w, h, color, line=None):
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(0.75)
        shp.shadow.inherit = False
        return shp

    def _text(self, s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=3):
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, Pt(1))
        first = True
        for item in runs:
            text, size, bold, color = item[0], item[1], item[2], item[3]
            italic = item[4] if len(item) > 4 else False
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(sa)
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = FONT
            r.font.color.rgb = color
        return tb

    def _brand(self, s, l, t, h=Inches(0.32)):
        if BRAND_LOGO and os.path.exists(BRAND_LOGO):
            s.shapes.add_picture(BRAND_LOGO, l, t, height=h)
        else:
            self._text(s, l, t, Inches(3), h, [(BRAND_TEXT, 13, True, RED)],
                       anchor=MSO_ANCHOR.MIDDLE)

    def _pic_fit(self, s, path, l, t, max_w, max_h, center="x"):
        pic = s.shapes.add_picture(path, l, t, width=max_w)
        if pic.height > max_h:
            sc = max_h / pic.height
            pic.height = int(pic.height * sc)
            pic.width = int(pic.width * sc)
        if center == "x":
            pic.left = int(l + (max_w - pic.width) / 2)
        elif center == "both":
            pic.left = int(l + (max_w - pic.width) / 2)
            pic.top = int(t + (max_h - pic.height) / 2)
        return pic

    # ----------------------------------------------------------------- chrome
    def _section(self, title, sub=None):
        s = self._slide()
        self.page += 1
        self._text(s, Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.7),
                   [(title, 30, False, RED)], anchor=MSO_ANCHOR.MIDDLE)
        if sub:
            self._text(s, Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.3),
                       [(sub, 12, False, GRAY)])
        self._brand(s, Inches(0.45), Inches(7.15))
        self._text(s, Inches(12.5), Inches(7.16), Inches(0.7), Inches(0.3),
                   [(str(self.page), 9, False, GRAY)], align=PP_ALIGN.RIGHT)
        return s

    def _footnote(self, s, text):
        self._text(s, Inches(0.5), Inches(6.44), Inches(12.3), Inches(0.62),
                   [(text, 8, False, GRAY, True)])

    def _bullets(self, s, x, y, w, bullets, size=12, spacing=0.52, head="Key takeaways"):
        if head:
            self._text(s, x, y, w, Inches(0.3), [(head + ":", size + 1.5, True, NAVY)])
            y = y + Inches(0.45)
        for b in bullets:
            sub = b.startswith("- ")
            txt = b[2:] if sub else b
            bx = x + (Inches(0.45) if sub else Inches(0.0))
            fs = size if not sub else size - 1.5
            boxw = w - Inches(0.22) - (bx - x)
            cpl = max(18, int((boxw / 914400) * 72 / (fs * 0.53)))
            nlines = max(1, math.ceil(len(txt) / cpl))
            self._rect(s, bx, y + Inches(0.07), Inches(0.09), Inches(0.09),
                       GOLD if not sub else GRAY)
            self._text(s, bx + Inches(0.22), y, boxw, Inches(0.26 * nlines + 0.1),
                       [(txt, fs, False, INK)])
            base = spacing if not sub else spacing - 0.06
            y += Inches(max(base, nlines * (fs * 1.2 / 72) + 0.16))
        return y

    def _dotgrid(self, s, x0, y0, cols, rows, step, color, d=Pt(5)):
        for r in range(rows):
            for c in range(cols):
                shp = s.shapes.add_shape(MSO_SHAPE.OVAL, int(x0 + c * step), int(y0 + r * step), d, d)
                shp.fill.solid(); shp.fill.fore_color.rgb = color
                shp.line.fill.background(); shp.shadow.inherit = False

    # ----------------------------------------------------------------- chart layouts
    def _notes(self, s, text):
        """Slides are charts only; the presenter brings their own notes, so we
        leave the speaker-notes pane empty (kept as a no-op hook)."""
        return

    def charts_bullets_slide(self, title, sub, imgs, notes=None, foot=None):
        """1-2 charts, no on-slide bullets; the talk-track lives in speaker notes."""
        s = self._section(title, sub)
        if len(imgs) == 1:
            self._pic_fit(s, imgs[0], Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.0),
                          center="both")
        else:
            self._pic_fit(s, imgs[0], Inches(0.4), Inches(1.5), Inches(6.16), Inches(4.7),
                          center="both")
            self._pic_fit(s, imgs[1], Inches(6.77), Inches(1.5), Inches(6.16), Inches(4.7),
                          center="both")
        self._notes(s, notes)
        if foot:
            self._footnote(s, foot)
        return s

    def wide_chart_slide(self, title, sub, img, notes=None, foot=None):
        """One large centered chart; the talk-track lives in speaker notes."""
        s = self._section(title, sub)
        self._pic_fit(s, img, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.0), center="both")
        self._notes(s, notes)
        if foot:
            self._footnote(s, foot)
        return s

    # ----------------------------------------------------------------- slides
    def title_slide(self):
        s = self._slide()
        self._rect(s, 0, 0, SW, SH, NAVY)
        self._rect(s, Inches(8.3), 0, SW - Inches(8.3), SH, RGBColor(0x24, 0x40, 0x70))
        self._text(s, Inches(0.7), Inches(0.5), Inches(6), Inches(0.4), [(self.asof, 14, False, LT_BLUE)])
        self._rect(s, Inches(0.72), Inches(1.04), Inches(6.6), Pt(1.2), RGBColor(0x4A, 0x66, 0x99))
        if BRAND_LOGO and os.path.exists(BRAND_LOGO):
            s.shapes.add_picture(BRAND_LOGO, Inches(10.5), Inches(0.45), height=Inches(0.5))
        else:
            self._text(s, Inches(9.0), Inches(0.45), Inches(3.8), Inches(0.5),
                       [(BRAND_TEXT, 18, True, WHITE)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        self._text(s, Inches(0.7), Inches(2.9), Inches(7.6), Inches(0.4),
                   [("Products Trading", 20, False, LT_BLUE)])
        self._text(s, Inches(0.68), Inches(3.4), Inches(8.0), Inches(1.7),
                   [("Refinery Outages", 38, True, WHITE),
                    (f"{FY} Outlook by Unit", 38, True, WHITE)], sa=2)
        self._rect(s, Inches(0.72), Inches(5.35), Inches(2.8), Pt(2.5), GOLD)
        self._text(s, Inches(0.72), Inches(5.55), Inches(7.9), Inches(0.9),
                   [("Capacity offline by unit, region & operator -- the gasoline and distillate read", 12.5, False, LT_BLUE),
                    (f"{FY}: ExxonMobil full-year (verified vs their plan); all other operators "
                     "H1-confirmed only", 11, False, RGBColor(0x9D, 0xB6, 0xDB))], sa=4)
        self._text(s, Inches(11.3), Inches(7.0), Inches(1.85), Inches(0.3),
                   [("PROPRIETARY", 9, False, LT_BLUE)], align=PP_ALIGN.RIGHT)

    def total_by_unit_slide(self):
        c2 = self.ctx["confirmed2027"]
        pk = lambda f: max(c2[f]["confirmed"])
        self.wide_chart_slide(
            f"Total {FY} Outages by Unit",
            "Capacity offline by unit & month; solid = confirmed, hatched = non-Exxon H2 (not yet booked)",
            self.a["splits_2027"],
            notes=(
                "2027 planned offline, per unit. CDU is the whole barrel: when crude is down, every "
                "product from the site is cut, gasoline and distillate alike. FCC and reformer feed the "
                "gasoline pool (cat gasoline and octane); the hydrocracker feeds distillate (diesel and "
                "jet). "
                f"Confirmed peak, all in H1: CDU ~{kbd(pk('CDU'))}, FCC ~{kbd(pk('FCC'))}, hydrocracker "
                f"~{kbd(pk('Hydrocracker'))}, reformer ~{kbd(pk('Reformer'))} kbd. "
                "Read each unit on its own, never added: a 250 CDU plus a 100 FCC is not '350 offline'. "
                "Solid bars are confirmed (Exxon full-year plus everyone's H1); the hatched autumn bars "
                "are non-Exxon H2, still being booked, so don't trade them as firm."),
            foot="Day-weighted offline (a unit down part of a month counts only its days down), each unit "
                 f"once per month. Non-Exxon H2 {FY} is a floor that fills in.")

    def drivers_slide(self):
        ev = engine.unit_events(self.ctx["df"], year=FY)
        ev = ev[ev["focus"].isin(engine.FOCUS_ORDER)].sort_values("kbd", ascending=False)
        top = ev.head(3)
        names = ", ".join(f"{r['plant'].replace(' Refinery', '')} ({kbd(r['kbd'])} kbd)"
                          for _, r in top.iterrows())
        mx = float(top.iloc[0]["kbd"]) if len(top) else 0.0
        p_tot = ev.head(12).groupby("padd")["kbd"].sum().sort_values(ascending=False)
        lead = p_tot.index[0] if len(p_tot) else "PADD 3"
        self.wide_chart_slide(
            "What's Driving the Numbers: the Biggest Outages",
            f"Each bar is one unit's nameplate offline (kbd) in {FY}, colored by PADD (region)",
            self.a["biggest_outages"],
            notes=(
                f"The biggest single units offline in 2027, by refinery and region. The most capacity "
                f"sits in {lead} (Gulf), so that's where one outage moves USGC supply and the export "
                "gasoline and distillate barrel most. Biggest single outages: "
                f"{names}. They're mostly crude (CDU): one crude unit down cuts the whole site, every "
                f"product. Read each bar on its own, never added: the biggest single outage is ~{kbd(mx)} "
                "kbd, not a summed total. Hatched bars are non-Exxon H2, still being booked, an "
                "indicative floor, not confirmed."),
            foot=f"Per-unit nameplate offline, the 12 biggest focus-unit outages of {FY}. Color = PADD "
                 "region; hatched = non-Exxon H2 (indicative).")

    def mom_slide(self):
        pc = self.ctx.get("period_change")
        if not pc:
            return                                              # graceful: need >= 2 reported months
        cur, prev = pc["cur_label"], pc["prev_label"]
        tc, tp = pc["total"]
        d = tc - tp
        arrow = "more" if d >= 0 else "less"
        bf = pc.get("by_focus")
        if bf is not None and len(bf):
            bf = bf[bf.index.isin(engine.FOCUS_ORDER)]
        lead_u = (bf.reindex(bf.abs().sort_values(ascending=False).index).index[0]
                  if (bf is not None and len(bf)) else "")
        bp = pc["by_padd"]
        lead_p = (str(bp.reindex(bp.abs().sort_values(ascending=False).index).index[0]).replace("PADD ", "P")
                  if len(bp) else "")
        self.wide_chart_slide(
            "Month-over-Month: What Changed and What's Driving It",
            f"Change in capacity offline, {cur} vs {prev}, by region (left) and focus unit (right)",
            self.a["mom_movers"],
            foot=(f"Day-weighted offline {arrow} month-on-month: ~{kbd(tc)} kbd in {cur} vs ~{kbd(tp)} in "
                  f"{prev} ({d:+,.0f} kbd). Biggest swing by region {lead_p}, by unit {lead_u}. Green = more "
                  "offline this month, red = less. Read each unit on its own."))

    def h1_compare_slide(self):
        h1 = self.ctx["h1_focus_planned"]

        def v(f, y):
            return float(h1.loc[f, y]) if (f in h1.index and y in h1.columns) else 0.0
        cdu27, cdu26, cdu25 = v("CDU", FY), v("CDU", FY - 1), v("CDU", FY - 2)
        d = (cdu27 / cdu26 - 1) if cdu26 else 0.0
        cm = self.ctx["focus_planned"]["CDU"]
        h1m = engine.MONTHS[:6]
        if FY in cm.index:
            pmo = max(h1m, key=lambda mo: cm.loc[FY, mo])
            pval = float(cm.loc[FY, pmo])
        else:
            pmo, pval = "Mar", 0.0
        self.wide_chart_slide(
            f"H1 Planned Outages per Unit & Month: {FY-2} / {FY-1} / {FY}",
            f"Like-for-like Jan-Jun planned offline, by unit and month ({FY} confirmed through H1)",
            self.a["h1_month_by_unit"],
            notes=(
                "H1 is the only honest cross-year read: 2027 is booked through June, H2 is still being "
                "scheduled. 2027 runs heavier than 2025 and 2026 in most H1 months (the orange bar). "
                f"Crude (CDU) averages ~{kbd(cdu27)} kbd vs ~{kbd(cdu26)} in 2026 and ~{kbd(cdu25)} in "
                f"2025, cresting in {pmo}. Crude and cat-cracker turnarounds cluster in February-March, "
                "landing into the spring gasoline-spec changeover, so the gasoline pool tightens just as "
                "summer-grade demand builds. The hydrocracker (distillate: diesel and jet) is front-"
                "loaded in Q1 too (Sweeny, Beaumont, Houston, LA, Martinez), but it is in line with "
                "history and below the 2023-24 peaks, not elevated. A heavier H1 = tighter gasoline and "
                "distillate supply into the spring."),
            foot="Day-weighted capacity offline by month, planned only, each unit once per month. "
                 f"Per-panel y-axis. {FY-2}/{str(FY-1)[2:]} actuals; {FY} the booked plan.")

    def gasoline_padd_slide(self):
        self.charts_bullets_slide(
            "Gasoline Complex: Crude (CDU) + Cat Cracker (FCC) by PADD",
            f"{FY} crude (CDU, left) and cat-cracker (FCC, right) offline by region & month -- the gasoline read",
            [self.a["cdu_padd_27"], self.a["fcc_padd_27"]],
            foot="Day-weighted concurrent offline by month, stacked by PADD. The FCC makes cat gasoline; a "
                 "CDU outage cuts the whole barrel, gasoline included. P3 (Gulf) is the swing region; P2 "
                 "(Midwest) carries the spring crude; P5 (West) is islanded. Past the dotted line = non-Exxon "
                 "H2 (a floor that grows as operators book).")

    def distillate_padd_slide(self):
        self.charts_bullets_slide(
            "Distillate Complex: Crude (CDU) + Hydrocracker by PADD",
            f"{FY} crude (CDU, left) and hydrocracker (right) offline by region & month -- the diesel/jet read",
            [self.a["cdu_padd_27"], self.a["hcu_padd_27"]],
            foot="Day-weighted concurrent offline by month, stacked by PADD. The hydrocracker makes diesel "
                 "and jet; the same CDU outage that cuts gasoline cuts distillate too (the whole barrel is "
                 "down). P3 (Gulf) and P2 (Midwest) carry the hydrocracker work. Past the dotted line = "
                 "non-Exxon H2 (indicative).")

    def naphtha_slide(self):
        nb = self.ctx["naphtha_balance"]
        net = nb["net"]
        order = sorted(range(12), key=lambda i: net[i])           # most negative first
        m1, m2 = engine.MONTHS[order[0]], engine.MONTHS[order[1]]
        v1, v2 = net[order[0]], net[order[1]]
        ny = int(round(nb["naphtha_yield"] * 100))
        state = "deficit" if nb["annual_net"] < 0 else "surplus"
        self.wide_chart_slide(
            "Naphtha Balance: CDU Supply vs Reformer Demand",
            f"{FY} outages read as naphtha length. CDU makes naphtha; reformers consume it",
            self.a["naphtha_balance"],
            notes=(
                f"Crude makes naphtha (~{ny}% of the barrel); reformers run on it to make reformate, the "
                "octane in gasoline. A CDU outage removes naphtha supply; a reformer outage removes "
                f"demand. All of 2027 sits in {state} (net {kbd(nb['annual_net'])} kbd): crude turnarounds "
                f"pull more naphtha off than reformer turnarounds free up, so naphtha is structurally "
                f"short, which is bullish reformate and octane. Tightest in {m1} ({kbd(v1)}) and {m2} "
                f"({kbd(v2)}) kbd, the autumn crude stack. When a reformer is down it frees naphtha but "
                "cuts ~85% of its feed as reformate, squeezing the gasoline pool's octane."),
            foot=f"Net = reformer offline x {nb['reformer_intake']:.0f} (demand) minus CDU offline x "
                 f"{nb['naphtha_yield']:.2f} (supply), day-weighted. + surplus / - deficit. "
                 "H2 non-Exxon is indicative.")

    def unplanned_context_slide(self):
        self.wide_chart_slide(
            f"Unplanned Offline: {FY-3}-{FY-1} Context",
            f"What unplanned actually looked like in recent years, to ground the {FY} scenario",
            self.a["unplanned_context"],
            notes=(
                "2027 has no unplanned actuals yet, so the scenario is anchored on what really happened. "
                "2024-2026 show the recurring shape: the February freeze spike and the September-October "
                "turnaround overlap, with summer quieter. This is the seasonal magnitude and range to "
                "carry into 2027 on top of the booked plan. Crude and the gasoline-making units (FCC, "
                "reformer) drive the spikes; distillate units add to the winter and fall windows. 2026 is "
                "reported through June, so its H2 is still filling in."),
            foot=f"Actual unplanned capacity offline by month, day-weighted. {FY-1} reported through June.")

    def scenario_slide(self):
        fan = self.ctx["scenario_fan"]
        pk_avg = float(max(fan["Average"]))
        self.wide_chart_slide(
            f"{FY} Unplanned Scenario: the Risk on Top of Planned",
            f"Potential unplanned offline (kbd per month) modeled on the {Y0}-{str(FY-1)[2:]} seasonal pattern",
            self.a["fan"],
            notes=(
                "2027 has no actual unplanned yet, so this is the monthly risk range to carry on top of "
                "the booked plan. Conservative, Average and Active scale the 2023-26 monthly shape by "
                "0.8, 1.0 and 1.3. Read it month by month, not as an annual total: risk peaks in February "
                "(freeze) and Sep-Oct (turnaround overlap), with summer the trough. The Average path "
                f"peaks near ~{kbd(pk_avg)} kbd of unplanned offline in the worst month. Trade the Active "
                "path as the supply-tightness stress case and Conservative as the floor."),
            foot=f"Scenario = mean {Y0}-{str(FY-1)[2:]} monthly unplanned shape (completeness-aware) x "
                 "{0.8 / 1.0 / 1.3}. Monthly concurrent offline, a risk range, not a forecast. Not an annual sum.")

    def market_slide(self):
        mc = engine.MARKET_CONTEXT
        src = mc["as_of"].split("(")[0].strip()
        self.wide_chart_slide(
            "What It Means for the Market: Gasoline & Distillate",
            "Spring crude turnarounds tighten both pools -- summer gasoline and a thin distillate cushion",
            self.a["market_setup"],
            foot=(f"Gasoline-complex (CDU + FCC + reformer) offline by month vs the summer-grade switchover. "
                  f"U.S. stocks below the 5-yr average on both sides: gasoline {mc['gasoline_vs_5yr_pct']:+d}%, "
                  f"distillate {mc['distillate_vs_5yr_pct']:+d}%, crude {mc['crude_vs_5yr_pct']:+d}% (EIA WPSR, "
                  f"{src}). Spring CDU turnarounds cut crude -- so gasoline and distillate both -- with the FCC "
                  "tightening summer gasoline and the hydrocracker tightening diesel/jet."))

    def build(self):
        self.title_slide()
        self.total_by_unit_slide()        # all four focus units, confirmed vs not-booked
        self.drivers_slide()              # biggest individual outages by PADD
        self.mom_slide()                  # month-over-month: what changed & what's driving it
        self.h1_compare_slide()           # H1 2025/26 vs 2027 planned, per unit & month
        self.gasoline_padd_slide()        # gasoline complex: CDU + FCC by PADD
        self.distillate_padd_slide()      # distillate complex: CDU + hydrocracker by PADD
        self.unplanned_context_slide()    # recent unplanned context
        self.scenario_slide()             # 2027 unplanned scenario (monthly)
        self.market_slide()               # what it means -- gasoline & distillate

    def save(self, path):
        self.prs.save(path)


def main():
    ap = argparse.ArgumentParser(description="Refinery outage trading-desk deck")
    ap.add_argument("excel", nargs="?", default=INPUT_PATH, help="path to the outage .xlsx export")
    ap.add_argument("--out", default=OUT_PATH, help="output .pptx path")
    args = ap.parse_args()

    print(f"Loading {args.excel.strip()} ...")
    ctx = engine.build_context(args.excel)
    with tempfile.TemporaryDirectory() as tmp:
        print("Rendering charts ...")
        assets = charts.render_all(ctx, tmp)
        print(f"Building deck -> {args.out}")
        Path(args.out).parent.mkdir(parents=True, exist_ok=True)
        deck = Deck(ctx, assets)
        deck.build()
        deck.save(args.out)
    print("Done.")


if __name__ == "__main__":
    sys.exit(main())
