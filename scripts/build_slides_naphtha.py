#!/usr/bin/env python3
"""
build_slides_naphtha.py
Parallel trading-desk deck, tilted to NAPHTHA / CHEM FEED and REFORMERS, and
headlining the in-progress year ("rest of <yy>") rather than the forward outlook
year. Same shape as the main deck (build_slides.py) - per-unit capacity offline,
chart-forward, presenter brings their own notes - but the read is octane and
petrochemical feedstock, not gasoline/distillate yields:

    1. Rest-of-year outages by unit   (H1 actual vs H2 booked, all four units)
    2. What's driving it              (biggest individual outages of the year)
    3. Reformers: the octane read     (reformer offline by month & PADD)   <- focus
    4. Naphtha balance                (CDU supply vs reformer demand, this year)
    5. Naphtha / octane / chem-feed   (the reforming-isom-aromatics complex)
    6. Reformer & crude by PADD       (where the octane / naphtha goes down)
    7. Unplanned context              (recent actuals, grounds the rest of year)

Naphtha is the reformer's charge and the steam-cracker / petrochemical feed, so
a reformer or naphtha-complex outage tightens octane and chem feed even when
crude runs hold - a read CDU-only trackers miss. Reuses the main deck's Deck
primitives and the shared engine context + charts, so the two decks agree.

Usage:
    python build_slides_naphtha.py                  # uses INPUT_PATH
    python build_slides_naphtha.py export.xlsx --out outage_deck_naphtha.pptx
"""
import argparse
import sys
import tempfile
from pathlib import Path

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import engine
import charts
from build_slides import (Deck, kbd, BRAND_LOGO, BRAND_TEXT,
                           NAVY, RED, GOLD, GREEN, LT_BLUE, GRAY, WHITE, SW, SH)
import os

_ROOT = Path(__file__).resolve().parent.parent
INPUT_PATH = str(_ROOT / "data" / "Golden_Record_Snowflake.xlsx")
OUT_PATH = str(_ROOT / "output" / "outage_deck_naphtha.pptx")

CY = engine.CURRENT_YEAR            # the in-progress year (rest of it) this deck headlines
FY = engine.FOCUS_YEAR             # plus the outlook year -- the forward window is CY + FY


class NaphthaDeck(Deck):
    """Reuses Deck's primitives/chrome; swaps in the naphtha/chem-feed slide set."""

    def title_slide(self):
        s = self._slide()
        self._rect(s, 0, 0, SW, SH, NAVY)
        self._rect(s, Inches(8.3), 0, SW - Inches(8.3), SH, RGBColor(0x24, 0x40, 0x70))
        self._text(s, Inches(0.7), Inches(0.5), Inches(6), Inches(0.4), [(self.asof, 14, False, LT_BLUE)])
        self._rect(s, Inches(0.72), Inches(1.04), Inches(6.6), Pt(1.2), RGBColor(0x4A, 0x66, 0x99))
        if BRAND_LOGO and os.path.exists(BRAND_LOGO):
            s.shapes.add_picture(BRAND_LOGO, Inches(10.5), Inches(0.45), height=Inches(0.5))
        else:
            self._text(s, Inches(9.0), Inches(0.45), Inches(3.8), Inches(0.5),
                       [(BRAND_TEXT, 18, True, WHITE)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        self._text(s, Inches(0.7), Inches(2.9), Inches(7.6), Inches(0.4),
                   [("Products Trading", 20, False, LT_BLUE)])
        self._text(s, Inches(0.68), Inches(3.4), Inches(8.0), Inches(1.7),
                   [("Refinery Outages", 38, True, WHITE),
                    (f"Rest of {str(CY)[2:]} + {str(FY)[2:]}: Naphtha, Reformers & Chem Feed", 27, True, WHITE)], sa=2)
        self._rect(s, Inches(0.72), Inches(5.35), Inches(2.8), Pt(2.5), GOLD)
        self._text(s, Inches(0.72), Inches(5.55), Inches(7.9), Inches(0.95),
                   [("The octane and petrochemical-feedstock read: reformers and the naphtha complex, "
                     "not just crude", 12.5, False, LT_BLUE),
                    (f"The forward window: {CY} H1 is actual; the rest of {CY} and all of {FY} are the outlook", 11,
                     False, RGBColor(0x9D, 0xB6, 0xDB))], sa=4)
        self._text(s, Inches(11.3), Inches(7.0), Inches(1.85), Inches(0.3),
                   [("PROPRIETARY", 9, False, LT_BLUE)], align=PP_ALIGN.RIGHT)

    def forward_unit_slide(self):
        self.wide_chart_slide(
            f"The Forward Book by Unit: Rest of {CY} + {FY}",
            f"Capacity offline by unit & month across the forward window ({CY} H1 shaded = actual)",
            self.a["focus_forward"],
            foot="Day-weighted offline, each unit once per month. Shaded = the current year's H1 (actual, behind "
                 "us); everything to the right is the forward book (rest of this year + next).")

    def reformer_slide(self):
        nbcy, nbfy = self.ctx["naphtha_balance_cy"], self.ctx["naphtha_balance"]
        fwd_ref = sum(nbcy["ref_offline"][6:]) + sum(nbfy["ref_offline"])
        fwd_reformate = sum(nbcy["reformate_lost"][6:]) + sum(nbfy["reformate_lost"])
        self.wide_chart_slide(
            "Reformers: the Octane Read",
            f"Catalytic-reformer offline by month & PADD across {CY}-{FY}, with reformate (octane) lost",
            self.a["reformer_forward"],
            foot=f"Reformers run naphtha -> reformate (the octane in gasoline). Over the forward window "
                 f"(rest of {CY} + {FY}), ~{kbd(fwd_ref)} kbd of reformer offline = ~{kbd(fwd_reformate)} kbd of "
                 "reformate (octane) not made. P3 (Gulf) carries the most.")

    def naphtha_balance_slide(self):
        nbcy, nbfy = self.ctx["naphtha_balance_cy"], self.ctx["naphtha_balance"]
        fwd_net = list(nbcy["net"][6:]) + list(nbfy["net"])
        n_def = sum(1 for v in fwd_net if v < -1e-6)
        ny = int(round(nbcy["naphtha_yield"] * 100))
        self.wide_chart_slide(
            "HVN Balance: CDU Supply vs Reformer Demand",
            f"Heavy virgin naphtha (the reformer feed) across {CY}-{FY}. CDU makes it; reformers consume it",
            self.a["naphtha_forward"],
            foot=f"Net = reformer offline (HVN demand) minus CDU offline x {nbcy['naphtha_yield']:.2f} (HVN supply, "
                 f"~{ny}% of crude), day-weighted. The forward window runs short ({n_def} of {len(fwd_net)} months "
                 f"in deficit): crude turnarounds pull more HVN off than reformers free.")

    def hvn_padd3_slide(self):
        p3cy, p3fy = self.ctx["naphtha_balance_cy_p3"], self.ctx["naphtha_balance_p3"]
        p3_net = list(p3cy["net"][6:]) + list(p3fy["net"])
        p3_def = sum(1 for v in p3_net if v < -1e-6)
        allnet = self.ctx["naphtha_balance"]["annual_net"]
        share = abs(p3fy["annual_net"]) / abs(allnet) if allnet else 0.0
        self.wide_chart_slide(
            "HVN Balance: PADD 3 (Gulf) Only",
            f"The same HVN balance, CDU & reformers filtered to PADD 3 (Gulf) -- where the tightness sits",
            self.a["naphtha_forward_p3"],
            foot=f"PADD 3 (Gulf) only -- the same supply/demand read as the prior slide, just the Gulf cut. "
                 f"{p3_def} of {len(p3_net)} forward months in deficit; the Gulf carries ~{share*100:.0f}% of the "
                 f"{FY} HVN deficit.")

    def chemfeed_slide(self):
        na = self.ctx["naphtha"]["annual"]
        cy_tot = float(na.loc[CY, "Total"]) if CY in na.index else 0.0
        fy_tot = float(na.loc[FY, "Total"]) if FY in na.index else 0.0
        self.wide_chart_slide(
            "Naphtha / Octane / Chem-Feed Complex",
            "Reforming, isomerization and aromatics/BTX offline by year - the octane & petrochem-feed read",
            self.a["naphtha_complex"],
            foot=f"The naphtha/octane complex (reformer charge + steam-cracker / petrochem feed). Offline "
                 f"~{kbd(cy_tot)} kbd in {CY} and ~{kbd(fy_tot)} kbd booked in {FY}. Reforming dominates; this "
                 "is the octane/chem-feed availability a CDU-only tracker misses.")

    def drivers_slide(self):
        self.charts_bullets_slide(
            f"What's Driving It: Biggest {CY} & {FY} Outages",
            f"Biggest individual focus-unit outages, {CY} (left) and {FY} (right), by PADD",
            [self.a["biggest_cy"], self.a["biggest_fy"]],
            foot="Per-unit nameplate offline (never summed), the biggest outages each year. Color = PADD region. "
                 f"{FY} non-Exxon H2 is still being booked (an indicative floor).")

    def padd_slide(self):
        self.charts_bullets_slide(
            "Reformer Outages by PADD",
            f"Reformer (octane) offline by region & month: {CY} (left) vs {FY} (right)",
            [self.a["ref_padd_cy"], self.a["ref_padd_fy"]],
            foot="Day-weighted concurrent reformer offline by month, stacked by PADD. P1 NE, P2 Midwest, "
                 "P3 Gulf, P4 Rockies, P5 West. P3 (Gulf) carries the bulk of the octane loss.")

    def unplanned_context_slide(self):
        self.wide_chart_slide(
            f"Unplanned Offline: {CY-2}-{CY} Context",
            f"What unplanned actually looked like recently, to ground the forward window",
            self.a["unplanned_context"],
            foot=f"Actual unplanned capacity offline by month, day-weighted. {CY} reported through June; "
                 "the Feb-freeze and autumn-turnaround windows are the recurring risk to carry forward.")

    def build(self):
        self.title_slide()
        self.forward_unit_slide()         # all four units across the forward window
        self.reformer_slide()             # reformers: the octane read  (forward, focus)
        self.naphtha_balance_slide()      # HVN balance, all PADDs (forward)
        self.hvn_padd3_slide()            # HVN balance, PADD 3 (Gulf) only -- one extra slide
        self.chemfeed_slide()             # naphtha/octane/chem-feed complex (annual)
        self.drivers_slide()              # biggest CY & FY outages
        self.padd_slide()                 # reformer by PADD, CY vs FY
        self.unplanned_context_slide()    # recent unplanned context


def main():
    ap = argparse.ArgumentParser(description="Naphtha / chem-feed rest-of-year deck")
    ap.add_argument("excel", nargs="?", default=INPUT_PATH, help="path to the outage .xlsx export")
    ap.add_argument("--out", default=OUT_PATH, help="output .pptx path")
    args = ap.parse_args()

    print(f"Loading {args.excel.strip()} ...")
    ctx = engine.build_context(args.excel)
    with tempfile.TemporaryDirectory() as tmp:
        print("Rendering naphtha/chem-feed charts ...")
        assets = charts.render_naphtha_assets(ctx, tmp)
        print(f"Building naphtha deck -> {args.out}")
        Path(args.out).parent.mkdir(parents=True, exist_ok=True)
        deck = NaphthaDeck(ctx, assets)
        deck.build()
        deck.save(args.out)
    print("Done.")


if __name__ == "__main__":
    sys.exit(main())
